import pickle
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import argparse

def extract_text_from_pptx(pptx_path):
    """
    Extract text content from each slide in the PowerPoint file
    """
    prs = Presentation(pptx_path)
    slides_content = []
    
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.strip())
        slides_content.append(' '.join(slide_text))
    
    # Remove empty slides
    slides_content = [text for text in slides_content if text.strip()]
    return slides_content

def load_transcript(pickle_path):
    """
    Load transcript from pickle file
    """
    with open(pickle_path, 'rb') as f:
        transcript = pickle.load(f)
    
    # If transcript is a dictionary or has specific format, modify accordingly
    if isinstance(transcript, dict):
        # Assuming transcript is stored under a 'text' key - modify as needed
        return transcript.get('text', '')
    return transcript

def evaluate_content_coverage(slides_content, transcript):
    """
    Evaluate content coverage using TF-IDF and cosine similarity
    """
    vectorizer = TfidfVectorizer(stop_words='english')
    
    # Prepare transcript segments
    words = transcript.split()
    chunk_size = max(len(words) // len(slides_content), 1)
    transcript_chunks = [' '.join(words[i:i + chunk_size]) 
                        for i in range(0, len(words), chunk_size)]
    
    # Transform texts to TF-IDF vectors
    all_text = slides_content + transcript_chunks
    tfidf_matrix = vectorizer.fit_transform(all_text)
    
    n_slides = len(slides_content)
    slides_vectors = tfidf_matrix[:n_slides]
    transcript_vectors = tfidf_matrix[n_slides:]
    
    similarity_matrix = cosine_similarity(slides_vectors, transcript_vectors)
    
    return {
        'average_similarity': float(np.mean(np.max(similarity_matrix, axis=1))),
        'min_similarity': float(np.min(np.max(similarity_matrix, axis=1))),
        'max_similarity': float(np.max(similarity_matrix))
    }

def evaluate_key_terms(slides_content, transcript):
    """
    Evaluate key terms coverage
    """
    vectorizer = TfidfVectorizer(stop_words='english')
    
    def extract_key_terms(text, n_terms=10):
        tfidf = vectorizer.fit_transform([text])
        feature_names = np.array(vectorizer.get_feature_names_out())
        sorted_idx = np.argsort(tfidf.toarray()[0])[::-1]
        return set(feature_names[sorted_idx][:n_terms])
    
    transcript_key_terms = extract_key_terms(transcript)
    slides_text = ' '.join(slides_content)
    slides_key_terms = extract_key_terms(slides_text)
    
    coverage = len(transcript_key_terms.intersection(slides_key_terms)) / len(transcript_key_terms)
    return float(coverage)

def calculate_information_density(slides_content):
    """
    Calculate information density metrics
    """
    def calc_density(text):
        words = text.split()
        return len(set(words)) / len(words) if words else 0
    
    densities = [calc_density(slide) for slide in slides_content]
    return {
        'average_density': float(np.mean(densities)),
        'density_variance': float(np.var(densities))
    }

def main():
    parser = argparse.ArgumentParser(description='Evaluate PowerPoint slides against video transcript')
    parser.add_argument('--pptx', required=True, help='Path to PowerPoint file')
    parser.add_argument('--transcript', required=True, help='Path to transcript pickle file')
    
    args = parser.parse_args()
    
    # Load and process files
    slides_content = extract_text_from_pptx(args.pptx)
    transcript = load_transcript(args.transcript)
    
    if not slides_content:
        print("Error: No text content found in slides")
        return
    
    # Run evaluations
    content_coverage = evaluate_content_coverage(slides_content, transcript)
    key_terms_coverage = evaluate_key_terms(slides_content, transcript)
    info_density = calculate_information_density(slides_content)
    
    # Print results
    print("\n=== Evaluation Results ===")
    print("\nContent Coverage:")
    print(f"Average Similarity: {content_coverage['average_similarity']:.3f}")
    print(f"Min Similarity: {content_coverage['min_similarity']:.3f}")
    print(f"Max Similarity: {content_coverage['max_similarity']:.3f}")
    
    print("\nKey Terms Coverage:")
    print(f"Coverage Score: {key_terms_coverage:.3f}")
    
    print("\nInformation Density:")
    print(f"Average Density: {info_density['average_density']:.3f}")
    print(f"Density Variance: {info_density['density_variance']:.3f}")
    
    # Print summary statistics about the slides
    print("\nSlide Statistics:")
    print(f"Number of slides with content: {len(slides_content)}")
    print(f"Average words per slide: {np.mean([len(slide.split()) for slide in slides_content]):.1f}")

if __name__ == "__main__":
    main()