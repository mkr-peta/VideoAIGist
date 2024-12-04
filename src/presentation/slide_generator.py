from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import Slide

class SlideGenerator:
    def __init__(self):
        self.prs = Presentation()
        # Set slide dimensions to standard 16:9
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self._setup_theme()
        
    def _setup_theme(self):
        """Setup presentation theme and layouts"""
        # Theme colors
        self.TITLE_COLOR = RGBColor(44, 62, 80)      # Dark blue
        self.SUBTITLE_COLOR = RGBColor(52, 73, 94)   # Slate
        self.ACCENT_COLOR = RGBColor(41, 128, 185)   # Blue
        self.TEXT_COLOR = RGBColor(44, 62, 80)       # Dark gray
        
        # Fonts
        self.TITLE_FONT = 'Calibri Light'
        self.BODY_FONT = 'Calibri'
        
        # Sizes
        self.TITLE_SIZE = Pt(54)  # Increased from 44
        self.SUBTITLE_SIZE = Pt(32)
        self.BODY_SIZE = Pt(18)
        self.NOTES_SIZE = Pt(12)

    def _add_background(self, slide: Slide):
        """Add a subtle gradient background"""
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = RGBColor(255, 255, 255)
        fill.gradient_stops[1].color.rgb = RGBColor(240, 244, 248)

    def _add_footer(self, slide: Slide, slide_number: int):
        """Add footer with slide number and decorative line"""
        # Add thin line at bottom
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            Inches(6.9),
            Inches(12.333),  # Adjusted for 16:9
            Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.ACCENT_COLOR
        line.line.fill.background()
        
        # Add slide number
        txBox = slide.shapes.add_textbox(
            Inches(12.333), Inches(6.8), Inches(0.5), Inches(0.3)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(slide_number)
        run.font.size = Pt(12)
        run.font.color.rgb = self.TEXT_COLOR

    def _set_text_properties(self, paragraph, font_size: Pt, font_name: str, 
                           color: RGBColor, bold: bool = False):
        """Helper to set text properties"""
        run = paragraph.add_run()
        run.font.size = font_size
        run.font.name = font_name
        run.font.color.rgb = color
        run.font.bold = bold
        return run

    def create_title_slide(self, title: str):
        """Create an attractive title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        self._add_background(slide)
        
        # Configure title box with proper positioning
        title_box = slide.shapes.title
        title_box.top = Inches(2.5)
        title_box.left = Inches(1.0)
        title_box.width = Inches(11.333)
        title_frame = title_box.text_frame
        title_frame.clear()
        
        # Set title text with proper formatting
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = self._set_text_properties(p, self.TITLE_SIZE, self.TITLE_FONT, 
                                      self.TITLE_COLOR, bold=True)
        run.text = title
        
        # Add decorative line centered below title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.67),
            Inches(4.2),
            Inches(2.0),
            Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.ACCENT_COLOR
        line.line.fill.background()

    def create_content_slide(self, content: dict, slide_number: int):
        """Create a beautifully formatted content slide with image support"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_background(slide)
        
        # Add title with adjusted spacing
        title_shape = slide.shapes.title
        title_shape.top = Inches(0.4)  # Keep this the same
        title_shape.left = Inches(0.5)
        title_shape.width = Inches(12.333)
        title_shape.height = Inches(1.0)  # Control title height
        title_shape.text = content['title']
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        title_frame.paragraphs[0].space_after = Pt(6)  # Reduce space after title
        run = title_frame.paragraphs[0].runs[0]
        run.font.size = self.SUBTITLE_SIZE
        run.font.name = self.TITLE_FONT
        run.font.color.rgb = self.TITLE_COLOR

        # Move content box up
        content_box = slide.placeholders[1]
        content_box.top = Inches(1.0)  # Changed from 1.5 to reduce gap
        
        # Check if slide has an image
        if 'image_path' in content:
            # More conservative width for content to prevent overlap
            content_width = Inches(6.8)  # Reduced from 7.5
            image_width = Inches(4.5)
            image_height = Inches(3.5)
            
            # Add image with more margin from text
            try:
                image = slide.shapes.add_picture(
                    content['image_path'],
                    left=Inches(8.3),  # Increased from 8.0 to move image right
                    top=Inches(1.8),
                    width=image_width,
                    height=image_height
                )
            except Exception as e:
                print(f"Failed to add image to slide {slide_number}: {str(e)}")
                content_width = Inches(11.933)  # Use full width if image fails
        else:
            content_width = Inches(11.933)  # Full width if no image
        
        # Configure content box with more precise positioning
        content_box = slide.placeholders[1]
        content_box.top = Inches(1)
        content_box.left = Inches(0.7)
        content_box.width = content_width
        tf = content_box.text_frame
        tf.clear()
        
        # Add bullet points with better spacing
        for point in content['points']:
            p = tf.add_paragraph()
            p.level = 0
            p.text = point
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(12)
            p.space_after = Pt(12)
            # Adjust line spacing
            p.line_spacing = 1.2  # 120% of font size
            for run in p.runs:
                run.font.size = self.BODY_SIZE
                run.font.name = self.BODY_FONT
                run.font.color.rgb = self.TEXT_COLOR
        
        # Add speaker notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.clear()
        
        for note in content['speaker_notes']:
            p = text_frame.add_paragraph()
            run = self._set_text_properties(p, self.NOTES_SIZE, self.BODY_FONT, 
                                        self.TEXT_COLOR)
            run.text = f"• {note}"
        
        self._add_footer(slide, slide_number)

    def generate_presentation(self, slides_content: list[dict], output_path: str, presentation_title: str = "Video Summary"):
        """Generate the complete presentation"""
        # Create title slide
        self.create_title_slide(presentation_title)
        
        # Create content slides
        for i, content in enumerate(slides_content, 1):
            self.create_content_slide(content, i)
            
        # Save presentation
        self.prs.save(output_path)
        print(f"Presentation saved to {output_path}")